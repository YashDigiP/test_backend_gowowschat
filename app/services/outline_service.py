# outline_service.py

from langchain_ollama import ChatOllama
from pptx import Presentation
from pptx.util import Inches
import tempfile
import os
from reportlab.lib.pagesizes import LETTER
from reportlab.pdfgen import canvas
import io
from services.llm_config import LLM_MODELS

def generate_outline(text, use_gemma=False):
    """
    Generate a professional PowerPoint outline based on input text.
    """
    model = LLM_MODELS["generate_insights"]
    print(f"📝 Using model for outline generation: {model}")
    llm = ChatOllama(model=model)
    prompt = f"""
You are an expert Presentation Designer.

Task:
- Read the below text carefully.
- Create a professional PowerPoint outline based on its content.

Rules:
- For each slide, give a main bullet (slide title) and 1-2 relevant sub-bullets (key points).
- Keep the titles concise and action-driven.
- No paragraph explanations. Only structured bullets.
- Limit total slides to 7 maximum unless critical.
- Maintain logical flow: Introduction -> Problem -> Solution -> Conclusion.

Input:
{text}

Output Example Format:

Slide 1: Introduction
- Background
- Objective

Slide 2: Challenges
- Problem 1
- Problem 2

Slide 3: Solutions
- Proposed Solution 1
- Benefits

(continue...)
"""

    try:
        result = llm.invoke(prompt)
        return result.content.strip()
    except Exception as e:
        return f"Error generating outline: {e}"


def export_outline_to_pptx(outline_text):
    """
    Convert outline text into a temporary PowerPoint file and return the path.
    """
    from pptx import Presentation
    prs = Presentation()
    slides = outline_text.strip().split("Slide ")[1:]

    for slide in slides:
        lines = slide.strip().split("\n")
        if not lines:
            continue
        title_line = lines[0].strip()
        content_lines = [line.strip("-").strip() for line in lines[1:] if line.strip()]
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide_obj = prs.slides.add_slide(slide_layout)
        slide_obj.shapes.title.text = f"Slide {title_line}"
        content_shape = slide_obj.placeholders[1]
        content_shape.text = "\n".join(content_lines)

    # ✅ Save to temporary file
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp_file.name)
    tmp_file.close()
    return tmp_file.name


def export_outline_to_pdf(outline_text: str):
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=LETTER)
    width, height = LETTER

    # Format and write the outline
    y = height - 40
    lines = outline_text.split('\n')

    for line in lines:
        if y < 40:
            pdf.showPage()
            y = height - 40

        pdf.drawString(40, y, line.strip())
        y -= 18

    pdf.save()
    buffer.seek(0)
    return buffer
